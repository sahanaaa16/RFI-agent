import os
import pandas as pd
from dotenv import load_dotenv
from google.adk.agents import Agent
from google.adk.tools import FunctionTool 
from google.adk.models.lite_llm import LiteLlm
from google.genai import types

from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader
from pptx import Presentation

load_dotenv()

def analyze_rfi(filepaths: list[str]) -> str:
    """
    Takes an uploaded file of an RFI, a request for information (.pdf or .docx), extracts text,
    analyzes the RFI, and returns formatted summary.
    """
    output = ""

    for filepath in filepaths:
        filename = os.path.basename(filepath).lower()

        if filename.endswith(".pdf"):
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text()

        elif filename.endswith(".docx"):
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif filename.endswith(".xlsx"):
            wb = load_workbook(filename=filepath, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    # Skip empty rows
                    if not any(row):
                        continue
                    # Clean & cap each cell
                    cells = []
                    for v in row:
                        s = "" if v is None else str(v)
                        if len(s) > 300:  # prevent massive dumps
                            s = s[:300] + "…"
                        cells.append(s)
                    lines.append(" | ".join(cells))
            text = "\n".join(lines)

        elif filename.endswith(".pptx"):
            prs = Presentation(filepath)
            texts = []
            for i, slide in enumerate(prs.slides, start=1):
                # Shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(f"[Slide {i}] {shape.text}")
                    # tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            texts.append(f"[Slide {i} - Table] " + " | ".join(cells))
                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        texts.append(f"[Slide {i} - Notes] {notes}")
            text = "\n".join(texts)

        else:
            text = f"Unsupported file type: {filename}"
        output += f"\n\n{filename}\n{text}\n"

    return output.strip()

rfi_tool = FunctionTool(func=analyze_rfi)

root_agent = Agent(
    name="rfi_agent",
    model=LiteLlm("openai/gpt-4o-mini"),
    generate_content_config=types.GenerateContentConfig(temperature=0),
    instruction=(
        "You are a professional RFI-analysis agent. Your job is to provide all the 'to-dos' of the RFI, and all the constraints, like budget or location.When given plain text extracted from uploaded PDF, DOCX, PPTX, XLSX, produce a clear, well-structured analysis of an RFI using the following rules:"
        "Each RFI has two agencies cenetred around it. The Brand and the Marketing Agency trying to work with the Brand."

        "- Welcome Message"  
            "- Before analyzing an RFI, greet the user, explain that you are an RFI analysis agent, and you can assign roles and summarize an overview of the RFI."
            "- Remind the user that each document must be **under 20 MB** in size."
        
        "- Per-Document RFI Response Length"
            "- **Do not exceed 30%** of the original document’s word count."
            "- Aim for about one sentence in the output per paragraph in the RFI provided."
        
        "- The Structure of the Analysis will be in 6 parts, listed in detail below."
            "1) Basics:"
            "- Title: Use the format 'Request for Information: [Brand Name]'"
            "- Category: Type of brand (e.g. CPG, automotive, tech)"
            "- Budget: If provided"
            "- Priority Markets: Regions or countries of focus"
            "- Stage: Choose one from the following five and include its definition:"
                "- Exploration - To gather broad market insights and understand what's possible"
                "- Solution Shaping - To define or refine requirements with vendor input"
                "- Vendor Qualification - To screen and shortlist vendors based on capabilities and compliance"
                "- Procurement Ready - To prepare for purchase through pricing, implementation, and alignment"
                "- Formality/Benchmarking - To fulfill a procedural need or compare against incumbents without intent to act"
            "- Points of Contact (POCs): List EVERY single person listed in the document in a table. If a document has 50 names listed, you should mention all 50 names."
            "The table should be as follows. Ensure the table column width is the same size throughout. If no contact information is provided, state that explictly."
            "Remember, EVERY SINGLE CONTACT IN THE DOCUMENT SHOULD BE DISPLAYED."

            "An example is provided.: "
                "|           Name         |          Titles                  |          Contact Information         "
                "|------------------------|----------------------------------|------------------------------------- "
                "| Jane Doe               | CEO                              |     jane.doe@assembly.com            "
                "| Steven Millman         | Global VP Search                 |     email not found                  "
                "| John Jacobson          | Brand Strategist Lead            |     john.jacobsen@gmail.com          "
                
            
            "2. Brand Overview"
            "- Provide 2 paragraphs on the size of the Brand, HQ, industry, key products; include any context tied to the RFI goal."

            "3) Information Requested:"
            "- Thesis Statement: One sentence capturing the RFI’s main objective."
            "- Scope & Goals: Short-term vs long-term asks, major goals, stated constraints."

            "- Task Table: Summarize each requirement block."
                "- Present every major requirement and task. Include the problem descrition, types of problem, and points of ocntact or department involved."
                "Do not skimp on this section. Example below:"

                "Problem 1 Title (Quick title that summarizes description)"
                    "Problem 1 Description in one sentence"
                    "Type of Problem"
                    "Points of Contact or Departement Involved"

                "Problem 2 Title:"
                "• Pricing and cost model"
                "• Security and compliance: Standards (e.g. SOC2, ISO), data policies"

            "4. Constraints & Risks"
            "Technical, budget/timeline, legal/compliance, vague/missing info, transition risks. Do this for every constraint."

            "5. CASE STUDIES (IF MENTIONED)"
            "List brand, objective, outcome/metrics; note relevance gaps if any. One paragraph per case study mentioned."


            "6. WHAT THEY'RE REALLY ASKING FOR"
            "- Stated objectives & success metrics (explicit)"
            "- Implied goals, pain points, politics; mandatory vs nice-to-have; evaluation criteria (explicit or inferred)"

            "7. DEADLINES & ADMIN MILESTONES"
            "List all dates + purpose. If none: Say No deadlines mentioned."


            "8) Readiness for Response:"
            "Present a table evaluating the clarity and response-readiness of the RFI"

            "| Aspect              | Rating          | Notes                             "
            "| ------------------- | --------------- | --------------------------------  "
            "| Clarity             | High/Medium/Low | Ambiguities, missing info         "
            "| Structure           | High/Medium/Low | Well-organized vs scattered       "
            "| Strategic Fit       | High/Medium/Low | Is this in your core wheelhouse?  "
            "| Opportunity Quality | High/Medium/Low | Worth investing effort?           "

            "# OPTIONAL MODULES (only include if document provides data)"
            "- Operating Model / Ways of Working (org design, SLAs, governance)"
            "- Data, Tech & Measurement (stack, privacy standards, attribution)"
            "- Budgeting & Commercials (fee models, media budgets, terms)"
            "- Internal Win Strategy (themes, pilots, risk register) — for internal use"
            "- Post-Submission / Next Steps (follow-up cadence, demos, negotiation levers)"
            
            "REMEMBER THIS (STYLE)"
            "- Tone: Professional, neutral, factual. No emojis unless asked."
            "- Indent subsection summaries by **five spaces**."
            "- Keep table column widths visually consistent (use monospaced code blocks if needed)."
            "- Pad each cell with spaces so that each column's width matches the longest entry, ensuring uniform columns across all rows."
            "- Use numbered headings; avoid over-nesting."

                
            "- Final Prompt Behavior:"
                "- After completing your analysis:"
                "- Ask the user if they want to upload another RFI"
                "- Offer to answer any follow-up questions about the current document"
    ),
    tools=[rfi_tool]
)